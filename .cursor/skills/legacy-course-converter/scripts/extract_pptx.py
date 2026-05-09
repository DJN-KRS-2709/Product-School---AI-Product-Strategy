#!/usr/bin/env python3
"""Extract a .pptx file to a single markdown file in insights/extracted/.

Usage:
    python extract_pptx.py <path-to-deck.pptx> [--out <output-dir>]

Output:
    insights/extracted/<basename>.md  (default; can be overridden with --out)

Each source slide becomes a `## Slide N — Title` section, with body text as
nested bullets, speaker notes as a quote block, and tables rendered as
markdown tables. Images, SmartArt, and embedded media are logged as
placeholders (the actual binary is not copied).
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    sys.stderr.write(
        "python-pptx not installed. Run: pip install -r "
        "requirements.txt (in this scripts/ folder)\n"
    )
    sys.exit(1)


def iter_text_runs(text_frame):
    for paragraph in text_frame.paragraphs:
        level = paragraph.level or 0
        text = "".join(run.text for run in paragraph.runs).strip()
        if text:
            yield level, text


def render_table(table) -> str:
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace("\n", " ") or " " for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
    if not rows:
        return ""
    header_sep = "| " + " | ".join(["---"] * len(table.columns)) + " |"
    return "\n".join([rows[0], header_sep, *rows[1:]])


def extract_slide(slide, slide_index: int) -> str:
    title = ""
    body_lines: list[str] = []
    table_blocks: list[str] = []
    image_logs: list[str] = []
    smartart_logs: list[str] = []
    media_logs: list[str] = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape == slide.shapes.title:
                title = shape.text_frame.text.strip().replace("\n", " ")
                continue
            for level, text in iter_text_runs(shape.text_frame):
                indent = "  " * level
                body_lines.append(f"{indent}- {text}")

        if shape.has_table:
            table_blocks.append(render_table(shape.table))

        shape_type = getattr(shape, "shape_type", None)
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_logs.append(
                f"[Image: {getattr(shape, 'name', 'picture')} — placement at slide {slide_index}]"
            )
        elif shape_type == MSO_SHAPE_TYPE.MEDIA:
            media_logs.append(f"[Media: {getattr(shape, 'name', 'media')} — slide {slide_index}]")
        elif shape_type == 24:
            smartart_logs.append(f"[SmartArt — manual review at slide {slide_index}]")

    notes = ""
    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame and notes_frame.text.strip():
            notes = notes_frame.text.strip()

    parts: list[str] = []
    heading = f"## Slide {slide_index}"
    if title:
        heading += f" — {title}"
    parts.append(heading)
    parts.append("")

    if body_lines:
        parts.extend(body_lines)
        parts.append("")

    for block in table_blocks:
        parts.append(block)
        parts.append("")

    for log in image_logs + smartart_logs + media_logs:
        parts.append(log)
    if image_logs or smartart_logs or media_logs:
        parts.append("")

    if notes:
        parts.append("> **Speaker notes:**")
        for line in notes.splitlines():
            parts.append(f"> {line}" if line.strip() else ">")
        parts.append("")

    return "\n".join(parts).rstrip() + "\n"


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("input", help="Path to the .pptx file")
    parser.add_argument(
        "--out",
        default=None,
        help="Output directory (default: insights/extracted/ relative to cwd)",
    )
    args = parser.parse_args()

    in_path = Path(args.input).expanduser().resolve()
    if not in_path.exists():
        sys.stderr.write(f"File not found: {in_path}\n")
        return 1

    out_dir = (
        Path(args.out).expanduser().resolve()
        if args.out
        else Path.cwd() / "insights" / "extracted"
    )
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"{in_path.stem}.md"

    prs = Presentation(str(in_path))

    sections: list[str] = [
        f"# {in_path.name}",
        "",
        f"Extracted from `{in_path}` — {len(prs.slides)} slides.",
        "",
        "---",
        "",
    ]

    for idx, slide in enumerate(prs.slides, start=1):
        sections.append(extract_slide(slide, idx))

    out_path.write_text("\n".join(sections), encoding="utf-8")
    print(f"Wrote {len(prs.slides)} slides to {out_path}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
